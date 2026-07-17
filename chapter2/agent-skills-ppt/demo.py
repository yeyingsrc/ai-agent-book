#!/usr/bin/env python3
"""
实验 2-6：使用 Agent Skills 从论文生成演示文稿（自建同构 Skills 机制）

本 demo 复现《深入理解 AI Agent》第二章「Agent Skills / 渐进式披露」一节的思想。
由于 Anthropic key 无效，这里用 OpenAI（gpt-4o-mini）+ 一套自建的、与 Anthropic
Skills 同构的机制来演示，核心是「渐进式披露（Progressive Disclosure）」：

  第一层（元数据）：Agent 启动时的 system prompt 里只放各 Skill 的 name +
                    description（薄目录，数百 token），并不含具体流程。
  第二层（核心流程）：当任务需要时，Agent 主动用 read_skill 工具加载完整 SKILL.md。
  第三层（细则）：Agent 可再用 read_skill_file 读取 reference.md / 脚本源码。

然后 Agent 用捆绑脚本 scripts/generate_pptx.py（经 run_skill_script 工具）用
python-pptx 生成真实的 .pptx，并读回校验页数与每页标题。

运行：
    export OPENAI_API_KEY=sk-...
    python demo.py
"""

import json
import os
import sys
from pathlib import Path

from openai import OpenAI
from pptx import Presentation

# 从同目录 .env 读取 OPENAI_API_KEY（若安装了 python-dotenv）
try:
    from dotenv import load_dotenv

    load_dotenv(Path(__file__).resolve().parent / ".env")
except ImportError:
    pass

# ---------------------------------------------------------------------------
# 路径与配置
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent
SKILLS_DIR = ROOT / "skills"
PAPER_PATH = ROOT / "papers" / "sample_paper.md"
OUTPUT_DIR = ROOT / "output"
MODEL = os.environ.get("OPENAI_MODEL", "gpt-4o-mini")


def log(msg: str) -> None:
    print(msg, flush=True)


# ---------------------------------------------------------------------------
# 第一层：启动时扫描 skills/ 目录，只读取每个 SKILL.md 的 frontmatter
# （name + description），拼成薄目录注入 system prompt。这一步刻意「只看目录」。
# ---------------------------------------------------------------------------
def parse_frontmatter(skill_md: str) -> dict:
    """从 SKILL.md 顶部的 --- YAML frontmatter --- 中解析 name / description。"""
    meta = {}
    if not skill_md.startswith("---"):
        return meta
    end = skill_md.find("---", 3)
    if end == -1:
        return meta
    for line in skill_md[3:end].splitlines():
        if ":" in line:
            k, v = line.split(":", 1)
            meta[k.strip()] = v.strip()
    return meta


def scan_skill_catalog() -> dict:
    """返回 {skill_name: {"description":..., "dir": Path}}，只含元数据。"""
    catalog = {}
    for skill_md in sorted(SKILLS_DIR.glob("*/SKILL.md")):
        meta = parse_frontmatter(skill_md.read_text(encoding="utf-8"))
        name = meta.get("name") or skill_md.parent.name
        catalog[name] = {
            "description": meta.get("description", ""),
            "dir": skill_md.parent,
        }
    return catalog


def build_system_prompt(catalog: dict) -> str:
    lines = [
        "你是一个能使用 Agent Skills 的助手。你并不预先知道每个 Skill 的详细流程，",
        "只在下方看到一份「薄目录」——每个 Skill 的 name 与 description（路由条件）。",
        "",
        "当任务需要某个 Skill 时，你必须：",
        "  1) 先用 read_skill(name) 加载它的完整 SKILL.md（第二层：核心流程）；",
        "  2) 如需实现/样式细节，再用 read_skill_file(name, path) 读取子文档或脚本（第三层）；",
        "  3) 按 SKILL.md 的约定，用 run_skill_script 调用捆绑脚本完成任务。",
        "不要在没有 read_skill 的情况下臆测某个 Skill 的调用方式。",
        "",
        "== 已安装的 Skills（薄目录，仅元数据）==",
    ]
    for name, info in catalog.items():
        lines.append(f"- {name}: {info['description']}")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# 工具实现：read_skill / read_skill_file / run_skill_script
# 这些是「渐进式披露」的通道——第二、三层内容只有被调用时才进入上下文。
# ---------------------------------------------------------------------------
def tool_read_skill(catalog: dict, name: str) -> str:
    info = catalog.get(name)
    if not info:
        return f"[error] 未找到 Skill: {name}"
    content = (info["dir"] / "SKILL.md").read_text(encoding="utf-8")
    log(f"\n  >>> [渐进式披露·第二层] Agent 调用 read_skill('{name}')，"
        f"加载完整 SKILL.md（{len(content)} 字符）")
    return content


def tool_read_skill_file(catalog: dict, name: str, rel_path: str) -> str:
    info = catalog.get(name)
    if not info:
        return f"[error] 未找到 Skill: {name}"
    target = (info["dir"] / rel_path).resolve()
    # 防目录穿越：必须落在该 skill 目录内
    if not str(target).startswith(str(info["dir"].resolve())):
        return f"[error] 非法路径: {rel_path}"
    if not target.exists():
        return f"[error] 文件不存在: {rel_path}"
    content = target.read_text(encoding="utf-8")
    log(f"  >>> [渐进式披露·第三层] Agent 调用 read_skill_file('{name}', '{rel_path}')，"
        f"加载子文档（{len(content)} 字符）")
    return content


def tool_run_skill_script(catalog: dict, name: str, script: str, payload: str) -> str:
    info = catalog.get(name)
    if not info:
        return f"[error] 未找到 Skill: {name}"
    script_path = (info["dir"] / "scripts" / script).resolve()
    if not script_path.exists():
        return f"[error] 脚本不存在: {script}"

    # 动态载入捆绑脚本（它就是 Skill 的一部分）
    import importlib.util
    spec = importlib.util.spec_from_file_location("bundled_generator", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    try:
        data = json.loads(payload) if isinstance(payload, str) else payload
    except json.JSONDecodeError as e:
        return f"[error] payload 不是合法 JSON: {e}"

    out_path = OUTPUT_DIR / "presentation.pptx"
    log(f"  >>> [执行捆绑脚本] run_skill_script('{name}', '{script}') "
        f"生成 {out_path.name} ...")
    result = module.build_presentation(data, str(out_path))
    return json.dumps(result, ensure_ascii=False)


TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "read_skill",
            "description": "加载指定 Skill 的完整 SKILL.md（核心流程，渐进式披露第二层）。",
            "parameters": {
                "type": "object",
                "properties": {"name": {"type": "string", "description": "Skill 名称"}},
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_skill_file",
            "description": "读取某 Skill 目录内的子文档或脚本源码（细则，渐进式披露第三层）。",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "path": {"type": "string", "description": "相对 skill 目录的路径，如 reference.md 或 scripts/generate_pptx.py"},
                },
                "required": ["name", "path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_skill_script",
            "description": "执行某 Skill 捆绑的脚本以完成实际产出（如生成 pptx）。",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "script": {"type": "string", "description": "脚本文件名，如 generate_pptx.py"},
                    "payload": {"type": "string", "description": "传给脚本的 JSON 字符串（大纲）"},
                },
                "required": ["name", "script", "payload"],
            },
        },
    },
]


def dispatch(catalog: dict, name: str, args: dict) -> str:
    if name == "read_skill":
        return tool_read_skill(catalog, args["name"])
    if name == "read_skill_file":
        return tool_read_skill_file(catalog, args["name"], args["path"])
    if name == "run_skill_script":
        return tool_run_skill_script(catalog, args["name"], args["script"], args["payload"])
    return f"[error] 未知工具: {name}"


# ---------------------------------------------------------------------------
# 主流程：agentic loop
# ---------------------------------------------------------------------------
def run_agent() -> Path | None:
    if not os.environ.get("OPENAI_API_KEY"):
        log("错误：未设置 OPENAI_API_KEY，请先 export OPENAI_API_KEY=sk-...")
        sys.exit(1)

    # timeout + 自动重试：单次网络/SSL 抖动不至于让整个 agentic loop 崩溃
    client = OpenAI(timeout=60.0, max_retries=3)
    catalog = scan_skill_catalog()

    system_prompt = build_system_prompt(catalog)
    log("=" * 72)
    log("【第一层·元数据】Agent 启动时只看到这份薄 Skill 目录（system prompt）：")
    log("-" * 72)
    log(system_prompt)
    log("-" * 72)
    log(f"（薄目录约 {len(system_prompt)} 字符 / 数百 token；各 Skill 的详细流程此刻并不在上下文中）")
    log("=" * 72)

    paper = PAPER_PATH.read_text(encoding="utf-8")
    user_task = (
        "请把下面这篇论文做成一份 8-12 页的演示文稿（含标题页、目录页、问题背景、"
        "方法概述、关键结果、局限性、小结页），总页数务必落在 8-12 页。"
        "先判断该用哪个 Skill，再严格按其 SKILL.md 的页序与约束操作。\n\n"
        "=== 论文全文 ===\n" + paper
    )

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_task},
    ]

    log("\n【任务下发】要求 Agent 从论文生成演示文稿。观察它如何按需渐进式披露：\n")

    final_result = None
    for turn in range(1, 9):
        resp = client.chat.completions.create(
            model=MODEL,
            messages=messages,
            tools=TOOLS,
            temperature=0.2,
        )
        msg = resp.choices[0].message
        messages.append(msg.model_dump(exclude_none=True))

        if not msg.tool_calls:
            log(f"\n【Agent 第 {turn} 轮·结束语】\n{msg.content}")
            break

        for tc in msg.tool_calls:
            fn = tc.function.name
            try:
                args = json.loads(tc.function.arguments or "{}")
            except json.JSONDecodeError:
                args = {}
            log(f"\n[Agent 第 {turn} 轮] 调用工具 -> {fn}({', '.join(f'{k}={_short(v)}' for k, v in args.items())})")
            result = dispatch(catalog, fn, args)
            if fn == "run_skill_script" and not result.startswith("[error]"):
                final_result = json.loads(result)
                log(f"  >>> 生成结果：{result}")
            messages.append({
                "role": "tool",
                "tool_call_id": tc.id,
                "content": result,
            })

    if final_result:
        return Path(final_result["path"])
    return None


def _short(v, n=48):
    s = str(v).replace("\n", " ")
    return s if len(s) <= n else s[:n] + "…"


# ---------------------------------------------------------------------------
# 校验：用 python-pptx 重新打开生成的文件，读回页数与每页标题，证明是有效 pptx。
# ---------------------------------------------------------------------------
def verify_pptx(path: Path) -> None:
    log("\n" + "=" * 72)
    log("【校验】用 python-pptx 重新打开生成的文件，读回页数与每页标题：")
    log("-" * 72)
    prs = Presentation(str(path))
    slides = list(prs.slides)
    log(f"文件: {path}")
    log(f"总页数: {len(slides)}")
    for i, slide in enumerate(slides, 1):
        first_text = "(空)"
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip():
                first_text = shp.text_frame.text.strip().splitlines()[0]
                break
        log(f"  第 {i:>2} 页标题: {first_text}")
    log("-" * 72)
    log(f"校验通过：这是一个可被 python-pptx / PowerPoint 打开的有效 .pptx（{len(slides)} 页）。")
    log("=" * 72)


def main():
    OUTPUT_DIR.mkdir(exist_ok=True)
    pptx_path = run_agent()
    if pptx_path and pptx_path.exists():
        verify_pptx(pptx_path)
    else:
        log("\n未生成 pptx。请检查上面的日志。")
        sys.exit(2)


if __name__ == "__main__":
    main()
